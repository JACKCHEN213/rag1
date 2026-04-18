"""
Document Parser Utility
Supports multiple document formats (PDF, Word, Excel, PPT, Markdown)
"""

import os
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import asyncio
import aiofiles
from io import BytesIO

# Optional imports for different document types
try:
    import fitz  # PyMuPDF
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class DocumentParser:
    """Universal document parser for multiple formats"""
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': self.parse_pdf,
            '.docx': self.parse_docx,
            '.doc': self.parse_doc,  # Note: .doc is harder to parse, may need additional libs
            '.xlsx': self.parse_xlsx,
            '.xls': self.parse_xlsx,
            '.pptx': self.parse_pptx,
            '.md': self.parse_markdown,
            '.markdown': self.parse_markdown,
            '.txt': self.parse_text,
            '.jpg': self.parse_image,
            '.jpeg': self.parse_image,
            '.png': self.parse_image,
        }
    
    async def parse_document(self, file_path: Path) -> List[Dict[str, Any]]:
        """
        Parse a document and return structured content elements
        
        Returns:
            List of elements, each with type, content, page_number, etc.
        """
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = file_path.suffix.lower()
        if ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {ext}")
        
        parse_func = self.supported_extensions[ext]
        return await parse_func(file_path)
    
    async def parse_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse PDF document"""
        if not PDF_AVAILABLE:
            raise ImportError("PyMuPDF not available. Install with: pip install PyMuPDF")
        
        elements = []
        
        # Open PDF
        doc = fitz.open(file_path)
        
        try:
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract text
                text = page.get_text()
                if text.strip():
                    elements.append({
                        "type": "text",
                        "content": text.strip(),
                        "page_number": page_num + 1,
                        "coordinates": None
                    })
                
                # Extract tables (simplified - in real implementation, use table detection)
                blocks = page.get_text("dict")["blocks"]
                for block in blocks:
                    if block.get("type") == 1:  # Image block
                        # This is a simplified check - real table detection needs more logic
                        block_text = block.get("lines", [])
                        if block_text:
                            elements.append({
                                "type": "table",
                                "content": str(block),  # Placeholder
                                "page_number": page_num + 1,
                                "coordinates": block.get("bbox")
                            })
                
                # Extract images
                imglist = page.get_images()
                for img_index, img in enumerate(imglist):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    
                    if pix.n > 4:  # CMYK: convert to RGB first
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    
                    # Save image temporarily or process OCR
                    img_data = pix.tobytes("png")
                    
                    elements.append({
                        "type": "image",
                        "content": f"<image_data_{page_num}_{img_index}>",  # Placeholder
                        "page_number": page_num + 1,
                        "metadata": {
                            "image_data": img_data,
                            "width": pix.width,
                            "height": pix.height
                        }
                    })
        finally:
            doc.close()
        
        return elements
    
    async def parse_docx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse Word document (.docx)"""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx not available. Install with: pip install python-docx")
        
        elements = []
        doc = Document(file_path)
        
        current_section = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Check for headings
                if paragraph.style.name.startswith('Heading'):
                    # Save previous section
                    if current_section:
                        elements.append({
                            "type": "text",
                            "content": "\n".join(current_section),
                            "page_number": None,
                            "metadata": {"style": "normal"}
                        })
                        current_section = []
                    
                    elements.append({
                        "type": "text",
                        "content": text,
                        "page_number": None,
                        "metadata": {"style": paragraph.style.name}
                    })
                else:
                    current_section.append(text)
        
        # Add remaining section
        if current_section:
            elements.append({
                "type": "text",
                "content": "\n\n".join(current_section),
                "page_number": None,
                "metadata": {"style": "normal"}
            })
        
        # Extract tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            elements.append({
                "type": "table",
                "content": str(table_data),  # Convert to string representation
                "page_number": None,
                "metadata": {"rows": len(table.rows), "cols": len(table.columns)}
            })
        
        return elements
    
    async def parse_doc(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse old Word format (.doc) - requires additional libraries"""
        # Note: .doc files are binary and harder to parse
        # For now, return a placeholder
        raise NotImplementedError(
            ".doc format parsing not yet implemented. "
            "Consider converting to .docx or using a library like antiword or textract."
        )
    
    async def parse_xlsx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse Excel document"""
        if not XLSX_AVAILABLE:
            raise ImportError("openpyxl not available. Install with: pip install openpyxl")
        
        elements = []
        wb = openpyxl.load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # Extract table data
            table_data = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    table_data.append([str(cell) if cell is not None else "" for cell in row])
            
            if table_data:
                elements.append({
                    "type": "table",
                    "content": str(table_data),
                    "page_number": None,  # Sheets don't have pages
                    "metadata": {
                        "sheet_name": sheet_name,
                        "rows": len(table_data),
                        "cols": max(len(row) for row in table_data) if table_data else 0
                    }
                })
        
        wb.close()
        return elements
    
    async def parse_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse PowerPoint document"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx not available. Install with: pip install python-pptx")
        
        elements = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            if slide_text:
                elements.append({
                    "type": "text",
                    "content": "\n\n".join(slide_text),
                    "page_number": slide_num + 1,
                    "metadata": {"slide_type": "content"}
                })
        
        return elements
    
    async def parse_markdown(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse Markdown document"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
        
        return await self.parse_markdown_content(content)
    
    async def parse_markdown_content(self, content: str) -> List[Dict[str, Any]]:
        """Parse markdown content"""
        elements = []
        lines = content.split('\n')
        
        current_section = []
        current_type = "text"
        
        for line in lines:
            line = line.rstrip()
            
            # Check for code blocks
            if line.startswith('```'):
                if current_type == "code":
                    # End of code block
                    if current_section:
                        elements.append({
                            "type": "code",
                            "content": "\n".join(current_section),
                            "page_number": None,
                            "metadata": {"language": "unknown"}
                        })
                        current_section = []
                    current_type = "text"
                else:
                    # Start of code block
                    if current_section:
                        elements.append({
                            "type": "text",
                            "content": "\n".join(current_section),
                            "page_number": None,
                            "metadata": {"style": "normal"}
                        })
                        current_section = []
                    current_type = "code"
                continue
            
            # Check for table rows (simplified)
            if '|' in line and current_type != "code":
                # This is a very basic table detection
                if current_type != "table" and current_section:
                    elements.append({
                        "type": "text",
                        "content": "\n".join(current_section),
                        "page_number": None,
                        "metadata": {"style": "normal"}
                    })
                    current_section = []
                current_type = "table"
                current_section.append(line)
                continue
            
            # Regular content
            if current_type in ["text", "table"] and line.strip():
                if current_type == "table":
                    # End of table
                    elements.append({
                        "type": "table",
                        "content": "\n".join(current_section),
                        "page_number": None,
                        "metadata": {"style": "markdown"}
                    })
                    current_section = [line]
                    current_type = "text"
                else:
                    current_section.append(line)
        
        # Add remaining content
        if current_section:
            if current_type == "code":
                elements.append({
                    "type": "code",
                    "content": "\n".join(current_section),
                    "page_number": None,
                    "metadata": {"language": "unknown"}
                })
            else:
                elements.append({
                    "type": "text",
                    "content": "\n".join(current_section),
                    "page_number": None,
                    "metadata": {"style": "normal"}
                })
        
        return elements
    
    async def parse_text(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse plain text document"""
        async with aiofiles.open(file_path, 'r', encoding='utf-8') as f:
            content = await f.read()
        
        return [{
            "type": "text",
            "content": content,
            "page_number": None,
            "metadata": {"style": "plain"}
        }]
    
    async def parse_image(self, file_path: Path) -> List[Dict[str, Any]]:
        """Parse image with OCR (placeholder)"""
        # Requires OCR library like pytesseract or PaddleOCR
        raise NotImplementedError(
            "Image parsing with OCR not yet implemented. "
            "Install an OCR library like PaddleOCR or pytesseract."
        )
    
    def extract_references(self, text: str) -> List[Dict[str, Any]]:
        """
        Extract figure/table references from text
        
        Patterns:
        - "Fig. 1", "Figure 1"
        - "Table 1", "Tab. 1"
        - "图1", "图-1"
        - "表1", "表-1"
        """
        references = []
        
        patterns = [
            # English patterns
            (r'Figure\s+\d+', 'figure'),
            (r'Fig\.\s*\d+', 'figure'),
            (r'Table\s+\d+', 'table'),
            (r'Tab\.\s*\d+', 'table'),
            # Chinese patterns
            (r'图\s*\d+', 'figure'),
            (r'图\s*[-\u2013]\s*\d+', 'figure'),
            (r'表\s*\d+', 'table'),
            (r'表\s*[-\u2013]\s*\d+', 'table'),
            # Combined patterns like "如图1所示"
            (r'如\s*图\s*\d+', 'figure'),
            (r'如\s*表\s*\d+', 'table'),
        ]
        
        for pattern, ref_type in patterns:
            matches = re.finditer(pattern, text, re.IGNORECASE)
            for match in matches:
                references.append({
                    "type": ref_type,
                    "match": match.group(0),
                    "start": match.start(),
                    "end": match.end(),
                    "number": re.search(r'\d+', match.group(0)).group(0) if re.search(r'\d+', match.group(0)) else None,
                    "context": text[max(0, match.start()-20):match.end()+20]
                })
        
        return references
    
    async def extract_text_from_image(self, image_path: Path) -> str:
        """Extract text from image using OCR"""
        # This would integrate with OCR library
        raise NotImplementedError("OCR integration not yet implemented")
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text"""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        # Remove null bytes
        text = text.replace('\x00', '')
        # Strip leading/trailing whitespace
        return text.strip()
